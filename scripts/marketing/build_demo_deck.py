#!/usr/bin/env python3
"""Build the Rytera live-demo PowerPoint (widescreen)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
PPT_PATH = ROOT / "marketing" / "assets" / "Rytera_Product_Demo.pptx"

NAVY = RGBColor(12, 15, 23)
PANEL = RGBColor(18, 24, 38)
CARD = RGBColor(24, 32, 50)
BRAND = RGBColor(91, 141, 239)
BRAND_LT = RGBColor(139, 176, 247)
CYAN = RGBColor(56, 189, 248)
VIOLET = RGBColor(167, 139, 250)
GREEN = RGBColor(52, 211, 153)
AMBER = RGBColor(251, 191, 36)
WHITE = RGBColor(248, 250, 252)
SOFT = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
RED = RGBColor(248, 113, 113)


def _set_run_font(run, *, size: int, color: RGBColor, bold: bool = False, name: str = "Calibri") -> None:
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def _textbox(slide, l, t, w, h, text, *, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, name="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=size, color=color, bold=bold, name=name)
    return box


def _para(tf, text, *, size=15, color=MUTED, bold=False, space_before=8, space_after=2):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=size, color=color, bold=bold)
    return p


def _rect(slide, l, t, w, h, fill: RGBColor, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    # tighter corners
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def _bar(slide, color=BRAND):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def _notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def _footer(slide, n: int, total: int) -> None:
    _textbox(slide, Inches(0.55), Inches(7.08), Inches(9), Inches(0.3), "Rytera™  ·  Confidential demo  ·  ryterainc.com", size=11, color=MUTED)
    _textbox(slide, Inches(11.6), Inches(7.08), Inches(1.4), Inches(0.3), f"{n}  /  {total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def _kicker(slide, text: str, y=0.32):
    _textbox(slide, Inches(0.55), Inches(y), Inches(12), Inches(0.32), text.upper(), size=12, color=BRAND_LT, bold=True)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 13

    # ── 1. Title ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s)
    _textbox(s, Inches(0.7), Inches(1.55), Inches(11), Inches(0.35), "PRODUCT DEMO  ·  12 MINUTES", size=13, color=BRAND_LT, bold=True)
    _textbox(s, Inches(0.7), Inches(2.05), Inches(12), Inches(1.3), "Stop hunting PDFs.\nStart underwriting.", size=40, color=WHITE, bold=True, name="Calibri")
    _textbox(
        s,
        Inches(0.7),
        Inches(4.55),
        Inches(11),
        Inches(0.9),
        "Rytera turns messy insurance, mortgage, and lending packages into a bind-ready memo —\nso licensed underwriters work the exceptions, not the stack.",
        size=18,
        color=SOFT,
    )
    _textbox(s, Inches(0.7), Inches(5.7), Inches(11), Inches(0.4), "Software  ·  B2B SaaS workbench  ·  Carriers, MGAs, line UW, staff UW, credit ops", size=14, color=MUTED)
    _footer(s, 1, TOTAL)
    _notes(
        s,
        """HOOK (first 60 seconds starts here — do not click around yet).
Open on this slide. Say: “Your desk was built for yes or no — not shared drives, re-key, and 11pm audit packs.”
Then: “In the next 12 minutes I will take one real-shaped file — Pacific Coast Distributors, a $4.35M warehouse — from broker package to a memo a licensed underwriter can sign. Bind stays off. Shadow stays on.”
Do not list features. Promise the journey.""",
    )

    # ── 2. What / who ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s)
    _kicker(s, "What you are watching")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), "A software product. Built for the desk that says yes or no.", size=28, color=WHITE, bold=True)

    cards = [
        (BRAND, "Product type", "B2B SaaS software — an AI underwriting workbench in the browser. Not hardware. Not a staffing service. Not a PAS replacement."),
        (CYAN, "Primary audience", "Chief underwriting officers, heads of UW, MGA principals, line UW managers at the branch, staff UW at the home office."),
        (VIOLET, "Also in the room", "Personal-lines leads, mortgage credit ops, commercial lenders, compliance / exam teams who will ask “why this file?”"),
        (GREEN, "How they buy", "Shadow pilot first. Bind off until their licensed UW cuts over. We do not invent a market appointment or a premium we cannot defend."),
    ]
    for i, (accent, title, body) in enumerate(cards):
        x = Inches(0.55 + (i % 2) * 6.3)
        y = Inches(1.65 + (i // 2) * 2.45)
        _rect(s, x, y, Inches(6.05), Inches(2.25), CARD)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), Inches(2.25))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        _textbox(s, x + Inches(0.35), y + Inches(0.22), Inches(5.5), Inches(0.4), title, size=16, color=accent, bold=True)
        _textbox(s, x + Inches(0.35), y + Inches(0.7), Inches(5.5), Inches(1.35), body, size=15, color=SOFT)
    _footer(s, 2, TOTAL)
    _notes(
        s,
        """Answer the room’s silent question in 20 seconds: this is software for underwriters and credit officers, not a robot that binds.
If a PAS vendor is in the room: “We sit in front of Guidewire / Duck Creek. We do not pretend bind is live until it is.”""",
    )

    # ── 3. Problem ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, RED)
    _kicker(s, "1. The hook  ·  The problem")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.9), "The file arrives messy. The decision still has to be defensible.", size=28, color=WHITE, bold=True)

    pains = [
        ("Days in the inbox", "Broker emails, ACORD XML, loss runs, SOVs, floor plans, W-2s — split across drives. Quote speed becomes a competitive wound."),
        ("Senior hours on hunting", "Licensed UW time goes to re-key and assembling audit packs, not appetite, terms, and price."),
        ("Black-box AI is a non-starter", "Examiners will ask “why this file?” A chatbot answer is not a SHA-256 audit ZIP."),
        ("PII into a model is a career risk", "Named insureds, SSNs, EINs, DOBs cannot go raw to an LLM API. Compliance will stop the project if they do."),
    ]
    for i, (h, b) in enumerate(pains):
        y = Inches(1.75 + i * 1.2)
        _rect(s, Inches(0.55), y, Inches(12.2), Inches(1.08), CARD)
        _textbox(s, Inches(0.8), y + Inches(0.12), Inches(11.7), Inches(0.35), h, size=16, color=WHITE, bold=True)
        _textbox(s, Inches(0.8), y + Inches(0.48), Inches(11.7), Inches(0.5), b, size=14, color=MUTED)
    _footer(s, 3, TOTAL)
    _notes(
        s,
        """Stay on pain. Mirror their words: shared drive, late submission, producer calling for a quote, exam next quarter.
Do not jump to the product yet. 20–25 seconds.""",
    )

    # ── 4. Value + glimpse ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, GREEN)
    _kicker(s, "1. The hook  ·  Value + glimpse of the result")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), "Underwriters work the exceptions. Rytera works the stack.", size=26, color=WHITE, bold=True)

    _rect(s, Inches(0.55), Inches(1.6), Inches(7.4), Inches(5.05), CARD)
    tf_box = s.shapes.add_textbox(Inches(0.85), Inches(1.85), Inches(6.9), Inches(4.6))
    tf = tf_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].clear()
    _para(tf, "What changes on the desk", size=14, color=BRAND_LT, bold=True, space_before=0)
    _para(tf, "Ingest the broker package once — Connect & pull, or drop files.", size=16, color=SOFT, space_before=12)
    _para(tf, "Extract with provenance. Conflicts surface as findings, not guesses.", size=16, color=SOFT)
    _para(tf, "Appetite, COPE, rating, memo — one visible journey.", size=16, color=SOFT)
    _para(tf, "A licensed underwriter still signs. Bind stays off until they say so.", size=16, color=SOFT)
    _para(tf, "Named insureds and PII are stripped before any LLM API call.", size=16, color=GREEN, bold=True)

    # Result card — the glimpse
    _rect(s, Inches(8.2), Inches(1.6), Inches(4.55), Inches(5.05), PANEL, line=BRAND)
    _textbox(s, Inches(8.45), Inches(1.8), Inches(4.15), Inches(0.3), "THE FILE WE WILL RUN", size=11, color=BRAND_LT, bold=True)
    _textbox(s, Inches(8.45), Inches(2.15), Inches(4.15), Inches(0.7), "Pacific Coast Distributors", size=20, color=WHITE, bold=True)
    facts = [
        ("Insured", "Warehouse · NAICS 493120"),
        ("TIV", "$4,350,000"),
        ("Line", "Property & business interruption"),
        ("Package", "ACORD · loss run · SOV · inspection"),
        ("Journey", "Ingest → Extract → Verify → Score → Price"),
        ("Human", "Sign-off still required"),
        ("Mode", "Shadow pilot · bind off"),
    ]
    for i, (k, v) in enumerate(facts):
        y = Inches(2.95 + i * 0.45)
        _textbox(s, Inches(8.45), y, Inches(1.35), Inches(0.4), k, size=12, color=MUTED)
        _textbox(s, Inches(9.85), y, Inches(2.7), Inches(0.4), v, size=13, color=SOFT, bold=True)
    _footer(s, 4, TOTAL)
    _notes(
        s,
        """Show the destination before the clicks. Read the Pacific Coast card aloud.
Then: “That is the golden path. I will not tour every menu. I will run this file.”
If live demo is ready, switch to the dashboard on the NEXT slide. If not, stay in the deck — screenshots / narrative still work.""",
    )

    # ── 5. Golden path ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s)
    _kicker(s, "2. The core workflow  ·  70% of the time")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.6), "One journey. The most common high-value file.", size=28, color=WHITE, bold=True)

    steps = [
        ("01", "Ingest", "Pull the broker package. No re-key of the stack."),
        ("02", "Extract", "ACORD, loss run, SOV, inspection — every field cited."),
        ("03", "Verify", "Conflicts and gaps become findings. Appetite fires first."),
        ("04", "Score", "Specialist agents + deterministic fallback. No black box."),
        ("05", "Price", "Indicated premium only when the book can defend it."),
        ("06", "Decide", "Memo + licensed UW checkpoint. Human still decides."),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = Inches(0.45 + i * 2.13)
        _rect(s, x, Inches(1.7), Inches(2.02), Inches(3.55), CARD)
        _textbox(s, x + Inches(0.12), Inches(1.9), Inches(1.78), Inches(0.4), num, size=18, color=BRAND_LT, bold=True)
        _textbox(s, x + Inches(0.12), Inches(2.4), Inches(1.78), Inches(0.7), title, size=18, color=WHITE, bold=True)
        _textbox(s, x + Inches(0.12), Inches(3.15), Inches(1.78), Inches(1.8), body, size=13, color=MUTED)
    _textbox(
        s,
        Inches(0.55),
        Inches(5.5),
        Inches(12.2),
        Inches(1.15),
        "Live click path: Dashboard → Insurance → Business / Commercial → Property & BI → Sample data “Pacific Coast” → Run.\n"
        "Narrate WHY at each step, not the button. Skip Settings, Registry, and every other vertical today.",
        size=15,
        color=SOFT,
    )
    _footer(s, 5, TOTAL)
    _notes(
        s,
        """GOLDEN PATH — do not deviate.
URL: /dashboard  (app.ryterainc.com/dashboard or staging).
Sign in as the demo underwriter. Theme: whatever is stable.
Click Insurance → Commercial → property_bi. Sample tab only if this product has a demo (Pacific Coast).
If Sample is hidden, use Files: drop the pacific-coast package.
Do not open mortgage or lending unless asked in Q&A.
Do not claim live CLUE / NCCI / Guidewire bind. Say “wired when your contracts are live; today this file runs on the honest pilot path.”""",
    )

    # ── 6. Ingest + extract ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, CYAN)
    _kicker(s, "Golden path  ·  Ingest & extract")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.6), "The stack becomes a package. Then every field has a source.", size=26, color=WHITE, bold=True)

    left = [
        ("Connect & pull", "SharePoint, Drive, S3, email, SFTP, IVANS-class intake — pull what the broker already sent. Why it matters: the UW never hunts the fourth PDF."),
        ("Drop files", "ACORD 125/126/130/140, loss run, SOV, inspection, floor plan, financials. Multi-file. Relevance scoring flags junk in the zip."),
        ("Writing company", "Pick the paper from the appointed panel. Rytera does not invent a market appointment you do not have."),
    ]
    right = [
        ("ACORD XML", "Named insured, locations, coverages, form numbers — structured data outranks a PDF guess."),
        ("Loss run & SOV", "Frequency, severity, schedule of values. Fake-clean loss history does not sneak onto the desk."),
        ("Floor plan / financials", "Area, stories, sprinklers; 16-line financial condition. COPE and creditworthiness without a second tool."),
    ]
    for i, (h, b) in enumerate(left):
        y = Inches(1.5 + i * 1.65)
        _rect(s, Inches(0.55), y, Inches(6.0), Inches(1.5), CARD)
        _textbox(s, Inches(0.8), y + Inches(0.15), Inches(5.5), Inches(0.35), h, size=16, color=CYAN, bold=True)
        _textbox(s, Inches(0.8), y + Inches(0.52), Inches(5.5), Inches(0.85), b, size=13, color=SOFT)
    for i, (h, b) in enumerate(right):
        y = Inches(1.5 + i * 1.65)
        _rect(s, Inches(6.8), y, Inches(6.0), Inches(1.5), CARD)
        _textbox(s, Inches(7.05), y + Inches(0.15), Inches(5.5), Inches(0.35), h, size=16, color=BRAND_LT, bold=True)
        _textbox(s, Inches(7.05), y + Inches(0.52), Inches(5.5), Inches(0.85), b, size=13, color=SOFT)
    _footer(s, 6, TOTAL)
    _notes(
        s,
        """WHILE CLICKING: hover the PII banner: “Named insureds and PII are stripped before any LLM API call — every insurance section, mortgage, and lending.”
That is the first wow. Pause two seconds. Let compliance people nod.
Then start the run. Watch the stage strip: Intake → Parse → Verify…""",
    )

    # ── 7. Verify score price ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, VIOLET)
    _kicker(s, "Golden path  ·  Verify, score, price, decide")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.6), "Findings, not vibes. A memo a human can sign.", size=26, color=WHITE, bold=True)

    items = [
        (GREEN, "Verify", "Cross-document conflicts (broker vs PDF vs SOV) surface as findings. Appetite fires before emotion. Out-of-appetite files refer."),
        (CYAN, "Score", "Specialist agents — risk, loss run, compliance, fraud, triage — each with a deterministic fallback. No LLM required to finish a file."),
        (AMBER, "Price", "ISO-style indicated premium when the loaded book can support it. If we cannot price it honestly, we do not invent a number."),
        (BRAND, "Decide", "ACCEPT / REFER / DECLINE memo. Conditions stay in sync with the evidence. Licensed UW license number on the sign-off."),
    ]
    for i, (c, h, b) in enumerate(items):
        x = Inches(0.55 + (i % 2) * 6.3)
        y = Inches(1.55 + (i // 2) * 2.5)
        _rect(s, x, y, Inches(6.05), Inches(2.3), CARD)
        _textbox(s, x + Inches(0.3), y + Inches(0.25), Inches(5.45), Inches(0.4), h, size=18, color=c, bold=True)
        _textbox(s, x + Inches(0.3), y + Inches(0.75), Inches(5.45), Inches(1.3), b, size=15, color=SOFT)
    _footer(s, 7, TOTAL)
    _notes(
        s,
        """Open the submission journey. Point at provenance: “This TIV came from the ACORD, not the model.”
Point at a finding: “This is why REFER exists — we would rather bounce a file than guess.”
Premium: if the demo book shows an indicated premium, say it is the pilot manual. Desk+ quotes their SERFF book. Never oversell bind.""",
    )

    # ── 8. Wow ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, GREEN)
    _kicker(s, "The wow  ·  What competitors usually skip")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), "The model never sees the named insured.", size=28, color=WHITE, bold=True)

    wows = [
        ("PII stripped first", "Named insureds, SSNs, EINs, dates of birth are redacted before any LLM API call — insurance, mortgage, and lending. That is the story compliance can say in an exam."),
        ("Every field is cited", "Structured broker data outranks AI-extracted PDF fields. Reconciliation is visible. No “the AI said so.”"),
        ("Zero-token first", "A router spends LLM budget only where rules genuinely fail. Most of the journey is deterministic. Cheap at scale. Works if the key is down."),
        ("Human still decides", "AI recommends. A licensed underwriter signs. Shadow first, bind last. Override analytics teach the book — they do not hide from the examiner."),
        ("Examiner-ready audit", "Encrypted at rest. SHA-256 manifest. Regulatory ZIP. When someone asks “why this file?” you hand them a bundle, not a war story."),
        ("Line desk + staff desk", "Branch coverage assist and producer service on one side. Home-office guides, rating plans, UW audits, and training on the other. Same platform."),
    ]
    for i, (h, b) in enumerate(wows):
        x = Inches(0.5 + (i % 3) * 4.2)
        y = Inches(1.55 + (i // 3) * 2.55)
        _rect(s, x, y, Inches(4.0), Inches(2.35), CARD)
        _textbox(s, x + Inches(0.22), y + Inches(0.2), Inches(3.55), Inches(0.55), h, size=16, color=GREEN if i == 0 else BRAND_LT, bold=True)
        _textbox(s, x + Inches(0.22), y + Inches(0.8), Inches(3.55), Inches(1.35), b, size=13, color=SOFT)
    _footer(s, 8, TOTAL)
    _notes(
        s,
        """THIS IS THE WOW. Slow down.
Click the audit trail / PII banner / a single cited field. One screenshot is worth ten adjectives.
Do not say “we are SOC 2 certified” unless that is true in the room’s materials. Say “SOC-ready audit trail” as the product feature.
Do not claim live carrier appointments.""",
    )

    # ── 9. Desks ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, CYAN)
    _kicker(s, "Same platform  ·  two desks")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.55), "Built the way the company is already organized.", size=26, color=WHITE, bold=True)

    _rect(s, Inches(0.55), Inches(1.5), Inches(6.05), Inches(5.1), CARD)
    _textbox(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(0.4), "Line UW desk  ·  branch", size=20, color=CYAN, bold=True)
    tf = s.shapes.add_textbox(Inches(0.85), Inches(2.25), Inches(5.45), Inches(4.1)).text_frame
    tf.word_wrap = True
    tf.paragraphs[0].clear()
    for t in (
        "Coverage assist on the live file",
        "Producer and policyholder service",
        "Quote / endorsement / certificate tickets",
        "Why it matters: the person facing the broker is not hunting PDFs while the producer is on the phone",
    ):
        _para(tf, "•  " + t, size=16, color=SOFT, space_before=10)

    _rect(s, Inches(6.8), Inches(1.5), Inches(6.05), Inches(5.1), CARD)
    _textbox(s, Inches(7.1), Inches(1.7), Inches(5.5), Inches(0.4), "Staff UW desk  ·  home office", size=20, color=VIOLET, bold=True)
    tf = s.shapes.add_textbox(Inches(7.1), Inches(2.25), Inches(5.45), Inches(4.1)).text_frame
    tf.word_wrap = True
    tf.paragraphs[0].clear()
    for t in (
        "UW guides and rating-plan reviews",
        "Market research and experience studies",
        "File audits and line-UW training",
        "Why it matters: the home office steers the book without another spreadsheet stack",
    ):
        _para(tf, "•  " + t, size=16, color=SOFT, space_before=10)
    _footer(s, 9, TOTAL)
    _notes(
        s,
        """Optional 45-second click: /line-uw then /staff-uw. If time is short, stay on this slide.
Also mention mortgage (Chen residential / Oak Street commercial) and lending (Blue Harbor Bakery) only if the audience is mixed credit + P&C.""",
    )

    # ── 10. Demo runbook ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, AMBER)
    _kicker(s, "3. Technical safeguards  ·  presenter runbook")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.55), "Do not demo production. Do not demo chaos.", size=26, color=WHITE, bold=True)

    rules = [
        ("Environment", "Use staging / the committed dashboard demo. Shadow on. Bind off. Notifications and chat pop-ups disabled. One browser profile. No extra tabs."),
        ("Offline backup", "Have a 90-second screen recording of Pacific Coast → memo ready. If the API hiccups, play it without apology and keep talking."),
        ("Clean UI", "Full screen. Hide bookmarks. Light or dark — pick one and stick. Zoom 110% so the back row can read the PII banner."),
        ("Honesty gates", "Do not claim live oracles, SERFF books, or Guidewire bind unless that tenant has them. Fail closed: no invented premium."),
        ("Data", "Pacific Coast Distributors, TIV $4.35M, warehouse. Never type “test” or “asdf”. Never show a real insured."),
        ("Timebox", "Hook 60s · journey ~8 min · wow 90s · wrap 60s. Q&A on the last slide with contact left up."),
    ]
    for i, (h, b) in enumerate(rules):
        y = Inches(1.4 + i * 0.85)
        _textbox(s, Inches(0.55), y, Inches(2.2), Inches(0.75), h, size=14, color=AMBER, bold=True)
        _textbox(s, Inches(2.8), y, Inches(9.9), Inches(0.8), b, size=14, color=SOFT)
    _footer(s, 10, TOTAL)
    _notes(
        s,
        """SKIP THIS SLIDE with the customer unless they ask how you run pilots. It is for your rehearsal.
Night-before: log in, run Pacific Coast once, leave the completed job in Recent Activity so the glimpse is instant.""",
    )

    # ── 11. Summary ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, GREEN)
    _kicker(s, "4. The wrap-up  ·  last 60 seconds")
    _textbox(s, Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), "What you just saw, in one breath.", size=28, color=WHITE, bold=True)

    recap = [
        ("Faster", "The queue ranks the book. The memo comes back in minutes, not days in an inbox."),
        ("Cheaper", "Senior UW hours go to yes or no — not hunting ACORDs and building audit packs."),
        ("Cleaner", "Appetite fires first. Findings are visible. Selection improves before loss ratio lectures you."),
        ("Defensible", "PII never hits the model raw. Encrypted audit. Licensed human still signs. Bind stays off until you cut over."),
    ]
    for i, (h, b) in enumerate(recap):
        x = Inches(0.55 + i * 3.15)
        _rect(s, x, Inches(1.7), Inches(3.0), Inches(4.4), CARD)
        _textbox(s, x + Inches(0.2), Inches(2.0), Inches(2.6), Inches(0.8), h, size=22, color=GREEN, bold=True)
        _textbox(s, x + Inches(0.2), Inches(2.9), Inches(2.6), Inches(2.8), b, size=15, color=SOFT)
    _footer(s, 11, TOTAL)
    _notes(
        s,
        """Read the four words. Then: “Underwriters work the exceptions. Rytera works the stack.”
Then go to CTA. Do not open another screen.""",
    )

    # ── 12. CTA ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s, BRAND)
    _kicker(s, "Call to action")
    _textbox(s, Inches(0.55), Inches(0.85), Inches(12), Inches(1.1), "30 minutes. Your book. Shadow on.\nBind off until your licensed UW says otherwise.", size=28, color=WHITE, bold=True)

    ctas = [
        ("1", "Book the deep-dive", "A walkthrough on your lines — commercial property, life, health, or credit — not a second marketing tour."),
        ("2", "Start a shadow pilot", "We measure accuracy on files you already decided. No PAS cutover. No invented bind."),
        ("3", "Email the team", "hello@ryterainc.com  ·  ryterainc.com  ·  dashboard at /dashboard"),
    ]
    for i, (n, h, b) in enumerate(ctas):
        y = Inches(2.35 + i * 1.35)
        _rect(s, Inches(0.55), y, Inches(12.2), Inches(1.2), CARD)
        _textbox(s, Inches(0.8), y + Inches(0.2), Inches(0.6), Inches(0.8), n, size=24, color=BRAND_LT, bold=True)
        _textbox(s, Inches(1.6), y + Inches(0.18), Inches(10.7), Inches(0.4), h, size=18, color=WHITE, bold=True)
        _textbox(s, Inches(1.6), y + Inches(0.6), Inches(10.7), Inches(0.45), b, size=14, color=MUTED)
    _footer(s, 12, TOTAL)
    _notes(
        s,
        """Ask for ONE next step: “Who owns the shadow-pilot decision, and can we lock 30 minutes this week on your top three commercial files?”
Do not offer a free-for-all sandbox without a named UW sponsor.""",
    )

    # ── 13. Q&A ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _bar(s)
    _textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4), "QUESTIONS", size=14, color=BRAND_LT, bold=True)
    _textbox(s, Inches(0.7), Inches(1.95), Inches(12), Inches(1.0), "We can stay on this slide.", size=40, color=WHITE, bold=True)
    _textbox(s, Inches(0.7), Inches(3.2), Inches(12), Inches(0.5), "Rytera Inc.  ·  AI underwriting workbench", size=18, color=SOFT)

    _rect(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.35), CARD)
    _textbox(s, Inches(1.0), Inches(4.35), Inches(11.2), Inches(0.4), "hello@ryterainc.com", size=22, color=WHITE, bold=True)
    _textbox(s, Inches(1.0), Inches(4.9), Inches(11.2), Inches(0.4), "ryterainc.com   ·   ryterainc.com/dashboard   ·   Book a demo on the site", size=16, color=SOFT)
    _textbox(s, Inches(1.0), Inches(5.5), Inches(11.2), Inches(0.6), "Likely Qs: live oracles · SERFF book · Guidewire bind · PII · line vs staff · mortgage/lending. Answer honestly: fail closed until those are actually live.", size=14, color=MUTED)
    _footer(s, 13, TOTAL)
    _notes(
        s,
        """Leave this up for the entire Q&A.
Honest answers:
- Live oracles: health-checked gateway; simulated on pilot; live when vendor contracts exist.
- Premium: demo book on pilot; their filed book on Desk+.
- Bind: off until licensed UW cutover. We do not re-key into a fake PAS.
- Appointments: they pick the writing company from their panel.
- Mortgage/lending: same PII strip, same journey, different docs (W-2, appraisal, credit).""",
    )

    PPT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_PATH)
    return PPT_PATH


if __name__ == "__main__":
    path = build()
    print(path)
