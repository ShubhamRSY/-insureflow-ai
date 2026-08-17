"""Shared data, theme, and slide helpers for C-suite marketing decks."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
ASSETS = ROOT / "marketing" / "assets"

# ── Theme ───────────────────────────────────────────────────────────────────
NAVY = RGBColor(12, 15, 23)
PANEL = RGBColor(18, 24, 38)
CARD = RGBColor(24, 32, 50)
BRAND = RGBColor(91, 141, 239)
BRAND_LT = RGBColor(139, 176, 247)
CYAN = RGBColor(56, 189, 248)
VIOLET = RGBColor(167, 139, 250)
GREEN = RGBColor(52, 211, 153)
AMBER = RGBColor(251, 191, 36)
WHITE = RGBColor(248, 250, 252)
SOFT = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
RED = RGBColor(248, 113, 113)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.55)
CONTENT_W = Inches(12.2)

BASELINE_MIN = 120
LOADED_UW_HR = 175
DESK_MO = 799
DESK_INCLUDED = 25
DESK_OVERAGE = 55
BOOK_MO = 2490
BOOK_INCLUDED = 80
BOOK_OVERAGE = 38


def _desk_annual_cost(memos_per_month: int) -> float:
    return (DESK_MO + max(0, memos_per_month - DESK_INCLUDED) * DESK_OVERAGE) * 12


def _book_annual_cost(memos_per_month: int) -> float:
    return (BOOK_MO + max(0, memos_per_month - BOOK_INCLUDED) * BOOK_OVERAGE) * 12


def _labor_return(memos_per_month: int, minutes_saved: float) -> float:
    hours = (minutes_saved / 60.0) * memos_per_month * 12
    return round(hours * LOADED_UW_HR, 0)


def _roi_pct(return_usd: float, cost_usd: float) -> float | None:
    if cost_usd <= 0:
        return None
    return round(((return_usd - cost_usd) / cost_usd) * 100.0, 0)


QUARTERLY = [
    {"q": "Q1", "phase": "Shadow pilot", "memos_mo": 20, "min_saved": 30, "pct_cycle_cut": 25, "software_mo": 350},
    {"q": "Q2", "phase": "Desk live", "memos_mo": 35, "min_saved": 60, "pct_cycle_cut": 50, "software_mo": DESK_MO + max(0, 35 - DESK_INCLUDED) * DESK_OVERAGE},
    {"q": "Q3", "phase": "Scale + STP", "memos_mo": 45, "min_saved": 75, "pct_cycle_cut": 62, "software_mo": DESK_MO + max(0, 45 - DESK_INCLUDED) * DESK_OVERAGE},
    {"q": "Q4", "phase": "Book tier", "memos_mo": 55, "min_saved": 90, "pct_cycle_cut": 75, "software_mo": BOOK_MO},
]

QUARTERLY_EXPLAIN = {
    "Q1": "Run 20 files in shadow. Bind stays off. UW compares AI memo to their decision. You learn override rate before you pay for live feeds.",
    "Q2": "Turn on Desk: live oracles + your SERFF rate book. First-pass time drops ~50%. This is the usual paid conversion point.",
    "Q3": "More volume, more straight-through Accepts. Senior UW hours shift to exceptions only.",
    "Q4": "Move to Book tier + Guidewire bind. Full payload, no re-key. Net benefit peaks.",
}

for row in QUARTERLY:
    ret_mo = (row["min_saved"] / 60.0) * row["memos_mo"] * LOADED_UW_HR
    row["labor_saved_mo"] = round(ret_mo, 0)
    row["net_benefit_mo"] = round(ret_mo - row["software_mo"], 0)
    row["roi_pct"] = _roi_pct(ret_mo * 12, row["software_mo"] * 12)

FY1_MEMOS = 40
FY1_MIN_SAVED = 75
FY1_RETURN = _labor_return(FY1_MEMOS, FY1_MIN_SAVED)
FY1_COST = _desk_annual_cost(FY1_MEMOS)
FY1_NET = FY1_RETURN - FY1_COST
FY1_ROI = _roi_pct(FY1_RETURN, FY1_COST)

FISCAL_YEARS = [
    {"fy": "FY1", "memos_mo": 40, "min_saved": 75, "cost_fn": _desk_annual_cost},
    {"fy": "FY2", "memos_mo": 65, "min_saved": 90, "cost_fn": _book_annual_cost},
    {"fy": "FY3", "memos_mo": 90, "min_saved": 95, "cost_fn": _book_annual_cost},
]
for fy in FISCAL_YEARS:
    fy["return_usd"] = _labor_return(fy["memos_mo"], fy["min_saved"])
    fy["cost_usd"] = fy["cost_fn"](fy["memos_mo"])
    fy["net_usd"] = fy["return_usd"] - fy["cost_usd"]
    fy["roi_pct"] = _roi_pct(fy["return_usd"], fy["cost_usd"])
    fy["pct_cycle_cut"] = round((fy["min_saved"] / BASELINE_MIN) * 100)

COMPANY_FY = [
    {"fy": "FY1", "customers": 8, "avg_mrr": 2200, "arr": 211_200, "gm_pct": 72},
    {"fy": "FY2", "customers": 22, "avg_mrr": 2800, "arr": 739_200, "gm_pct": 78},
    {"fy": "FY3", "customers": 48, "avg_mrr": 3400, "arr": 1_958_400, "gm_pct": 82},
]

BASELINE_LABOR_MO = round((BASELINE_MIN / 60) * 40 * LOADED_UW_HR, 0)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    return s


def _set_run(run, *, size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


def textbox(slide, left, top, width, height, text, *, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, color=color, bold=bold)
    return box


def rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def bar(slide, color: RGBColor = BRAND) -> None:
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), SLIDE_H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def footer(slide, n: int, total: int, label: str = "Confidential · Planning model labeled") -> None:
    textbox(slide, MARGIN_L, Inches(7.06), Inches(9.5), Inches(0.28), label, size=10, color=MUTED)
    textbox(slide, Inches(11.65), Inches(7.06), Inches(1.35), Inches(0.28), f"{n} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def slide_header(
    slide,
    *,
    kicker: str,
    title: str,
    subtitle: str = "",
    explain: str = "",
    accent: RGBColor = BRAND,
) -> float:
    """Return Y position (inches) where body content should start."""
    bar(slide, accent)
    textbox(slide, MARGIN_L, Inches(0.28), CONTENT_W, Inches(0.28), kicker.upper(), size=11, color=BRAND_LT, bold=True)
    textbox(slide, MARGIN_L, Inches(0.58), CONTENT_W, Inches(0.55), title, size=24, color=WHITE, bold=True)
    y = 1.18
    if subtitle:
        textbox(slide, MARGIN_L, Inches(y), CONTENT_W, Inches(0.38), subtitle, size=14, color=SOFT)
        y += 0.42
    if explain:
        rect(slide, MARGIN_L, Inches(y), CONTENT_W, Inches(0.72), PANEL, line=BRAND)
        textbox(slide, Inches(0.75), Inches(y + 0.08), Inches(11.5), Inches(0.58), explain, size=12, color=MUTED)
        y += 0.82
    return y


def add_table(
    slide,
    left,
    top,
    width,
    rows: list[list[str]],
    *,
    col_fracs: list[float] | None = None,
    header: bool = True,
    row_height: float = 0.42,
    font_size: int = 11,
    right_cols: set[int] | None = None,
) -> float:
    """Native pptx table. col_fracs must sum to 1.0. Returns bottom Y in inches."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    if col_fracs is None:
        col_fracs = [1.0 / n_cols] * n_cols
    total_frac = sum(col_fracs)
    col_fracs = [f / total_frac for f in col_fracs]

    height = Inches(row_height * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, Inches(top), width, height)
    table = shape.table

    col_widths = [int(width * f) for f in col_fracs]
    # fix rounding drift on last column
    col_widths[-1] = int(width) - sum(col_widths[:-1])
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    right_cols = right_cols or set()

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if ci in right_cols else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text

            if ri == 0 and header:
                _set_run(run, size=font_size, color=BRAND_LT, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL
            else:
                is_roi = "%" in cell_text and ci == n_cols - 1
                _set_run(run, size=font_size, color=GREEN if is_roi else SOFT, bold=is_roi)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD if ri % 2 else PANEL

    return top + row_height * n_rows


def explain_card(slide, left, top, width, height, title: str, body: str, accent: RGBColor = BRAND_LT):
    rect(slide, left, top, width, height, CARD)
    textbox(slide, left + Inches(0.22), top + Inches(0.14), width - Inches(0.35), Inches(0.32), title, size=13, color=accent, bold=True)
    textbox(slide, left + Inches(0.22), top + Inches(0.48), width - Inches(0.35), height - Inches(0.55), body, size=11, color=SOFT)
