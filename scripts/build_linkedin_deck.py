#!/usr/bin/env python3
"""Capture Rytera dashboard screenshots and build a professional product deck."""

from __future__ import annotations

import json
import sys
import urllib.request
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "assets" / "linkedin_deck"
PPT_PATH = ROOT / "Rytera_LinkedIn_Deck.pptx"
BASE = "http://127.0.0.1:8002"
DASHBOARD = f"{BASE}/dashboard"
JOB_ID = "demo-linkedin-001"
LOGIN = {"username": "demo_uw", "password": "demopass123"}


def api_login() -> str:
    req = urllib.request.Request(
        f"{BASE}/auth/login",
        data=json.dumps(LOGIN).encode(),
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=15) as resp:
        return json.loads(resp.read())["access_token"]


def capture_screenshots(token: str) -> dict[str, Path]:
    from playwright.sync_api import sync_playwright

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    shots: dict[str, Path] = {}

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(viewport={"width": 1440, "height": 900}, device_scale_factor=2)
        page = context.new_page()

        page.goto(DASHBOARD, wait_until="domcontentloaded")
        page.evaluate(
            """([token, user]) => {
                localStorage.setItem('insureflow_token', token);
                localStorage.setItem('insureflow_user', JSON.stringify(user));
            }""",
            [token, {"username": LOGIN["username"], "role": "underwriter", "org_id": "acme"}],
        )
        page.reload(wait_until="networkidle")
        page.wait_for_timeout(1500)

        def snap(name: str) -> None:
            path = OUT_DIR / f"{name}.png"
            page.screenshot(path=str(path), full_page=False)
            shots[name] = path

        snap("01_overview")
        page.goto(f"{DASHBOARD}/insurance", wait_until="networkidle")
        page.wait_for_timeout(2000)
        snap("02_insurance")

        row = page.locator("tr").filter(has_text=JOB_ID).first
        if row.count():
            row.click()
            page.wait_for_timeout(2500)
            snap("03_submission_journey")

        page.goto(f"{DASHBOARD}/queue", wait_until="networkidle")
        page.wait_for_timeout(1500)
        snap("04_queue")
        browser.close()

    return shots


def build_ppt(shots: dict[str, Path]) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK = RGBColor(12, 15, 23)
    PANEL = RGBColor(18, 24, 38)
    BRAND = RGBColor(99, 102, 241)
    BRAND_LIGHT = RGBColor(129, 140, 248)
    MUTED = RGBColor(148, 163, 184)
    WHITE = RGBColor(241, 245, 249)
    SOFT = RGBColor(203, 213, 225)

    slide_num = 0

    def _dark_bg(slide) -> None:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK

    def _accent_bar(slide) -> None:
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BRAND
        bar.line.fill.background()

    def _footer(slide, label: str = "Rytera™ · rytera.ai") -> None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED

    def _slide_number(slide) -> None:
        nonlocal slide_num
        slide_num += 1
        box = slide.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = str(slide_num)
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.RIGHT

    def _add_bullets(tf, items: list[str], *, size: int = 15, color=MUTED, bold_first: bool = False) -> None:
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.level = 0
            p.space_before = Pt(8 if i else 0)
            p.space_after = Pt(4)
            if bold_first and i == 0:
                p.font.bold = True
                p.font.color.rgb = SOFT

    def add_title_slide() -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)
        _accent_bar(slide)

        tag = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4), Inches(0.4))
        tag.text_frame.paragraphs[0].text = "PRODUCT OVERVIEW"
        tag.text_frame.paragraphs[0].font.size = Pt(11)
        tag.text_frame.paragraphs[0].font.color.rgb = BRAND_LIGHT
        tag.text_frame.paragraphs[0].font.bold = True

        title = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(11.5), Inches(1.2))
        tp = title.text_frame.paragraphs[0]
        tp.text = "Rytera™"
        tp.font.size = Pt(54)
        tp.font.bold = True
        tp.font.color.rgb = WHITE

        sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.35), Inches(11), Inches(1.2))
        sp = sub.text_frame.paragraphs[0]
        sp.text = "AI underwriting with full pipeline visibility"
        sp.font.size = Pt(24)
        sp.font.color.rgb = SOFT

        desc = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(10.5), Inches(1.5))
        _add_bullets(
            desc.text_frame,
            [
                "Commercial insurance · Mortgage · Lending",
                "From document intake to bind-ready decision — every layer explained, not hidden",
            ],
            size=16,
            color=MUTED,
        )
        _footer(slide, "Rytera™ · rytera.ai · Product overview deck")
        _slide_number(slide)

    def add_text_slide(heading: str, lead: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)
        _accent_bar(slide)

        h = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12), Inches(0.8))
        h.text_frame.paragraphs[0].text = heading
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = WHITE

        lead_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.5), Inches(0.9))
        lp = lead_box.text_frame.paragraphs[0]
        lp.text = lead
        lp.font.size = Pt(17)
        lp.font.color.rgb = SOFT

        body = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11.8), Inches(4.2))
        _add_bullets(body.text_frame, bullets, size=16)

        _footer(slide)
        _slide_number(slide)

    def add_feature_slide(
        heading: str,
        what_it_is: str,
        bullets: list[str],
        image: Path,
        caption: str,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)
        _accent_bar(slide)

        # Left narrative panel
        panel = slide.shapes.add_shape(1, Inches(0.45), Inches(0.45), Inches(4.85), Inches(6.35))
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL
        panel.line.color.rgb = RGBColor(40, 48, 68)
        panel.line.width = Pt(0.75)

        h = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(4.4), Inches(0.9))
        h.text_frame.word_wrap = True
        hp = h.text_frame.paragraphs[0]
        hp.text = heading
        hp.font.size = Pt(22)
        hp.font.bold = True
        hp.font.color.rgb = WHITE

        intro = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(4.4), Inches(1.1))
        intro.text_frame.word_wrap = True
        ip = intro.text_frame.paragraphs[0]
        ip.text = what_it_is
        ip.font.size = Pt(13)
        ip.font.color.rgb = BRAND_LIGHT

        body = slide.shapes.add_textbox(Inches(0.7), Inches(2.85), Inches(4.4), Inches(3.5))
        body.text_frame.word_wrap = True
        _add_bullets(body.text_frame, bullets, size=12, color=MUTED)

        # Right screenshot
        slide.shapes.add_picture(str(image), Inches(5.55), Inches(0.55), width=Inches(7.35))

        cap = slide.shapes.add_textbox(Inches(5.55), Inches(6.55), Inches(7.35), Inches(0.4))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(10)
        cp.font.color.rgb = MUTED
        cp.alignment = PP_ALIGN.CENTER

        _footer(slide)
        _slide_number(slide)

    # ── Deck content ────────────────────────────────────────────────

    add_title_slide()

    add_text_slide(
        "The problem",
        "Most underwriting tools treat AI like a black box.",
        [
            "Brokers upload a package; the system returns accept, refer, or decline.",
            "Underwriters cannot see how documents were parsed, verified, or priced.",
            "Compliance and audit teams lack a clear trail of agent reasoning and data sources.",
            "Carriers need transparency — not just speed — before trusting AI in production.",
        ],
    )

    add_text_slide(
        "What Rytera does",
        "Rytera is an AI underwriting operating system for insurance, mortgage, and lending.",
        [
            "Ingests multi-format submissions: ACORD, loss runs, tax returns, credit, appraisals, and more.",
            "Runs a structured pipeline: intake → parse → verify → score → price → decide.",
            "Surfaces every step in the dashboard so underwriters review the work, not just the outcome.",
            "Connects to carrier systems — oracles, policy admin, CRM, loss control — via an integration gateway.",
            "Produces memos, quotes, encrypted audit trails, and licensed underwriter sign-off.",
        ],
    )

    feature_slides = [
        (
            "01_overview",
            "Command center",
            "The Overview page is the home screen for underwriting operations. It shows market context, job volume, and the end-to-end pipeline at a glance.",
            [
                "Pipeline strip — maps each submission through Intake, Parse, Verify, Score, Price, and Decide.",
                "Three verticals — Commercial P&C, Mortgage, and Lending on one shared platform.",
                "Job metrics — live counts for insurance and mortgage workloads plus pending UW sign-offs.",
                "Quick demos — one-click sample submissions for sales and training.",
                "Recent activity — open any job to inspect its full submission journey.",
            ],
            "Screenshot: Overview dashboard with pipeline stages and vertical coverage",
        ),
        (
            "02_insurance",
            "Commercial insurance workflow",
            "The Insurance workspace manages P&C submissions from broker intake through bind-ready output.",
            [
                "Underwriting pipeline — visual flow from document intake to final decision.",
                "Integration feeds — health status for CLUE, NCCI, Guidewire, loss control, and related systems.",
                "Source hub — pull packages from connected folders, APIs, or demo connectors.",
                "Job table — every submission shows status, decision, premium, and journey progress.",
                "Demo presets — run realistic scenarios such as Pacific Coast Marine instantly.",
            ],
            "Screenshot: Commercial insurance page with pipeline and job list",
        ),
        (
            "03_submission_journey",
            "Submission journey panel",
            "Click any job to open the Submission Journey — the core differentiator of Rytera.",
            [
                "Timeline — stage-by-stage progress with duration and status for each step.",
                "COPE & risk — construction, occupancy, protection, and exposure analysis.",
                "Oracle verification — external data lookups (CLUE, NCCI, A-PLUS, catastrophe models).",
                "Agent findings — specialist AI agents flag risk, compliance, fraud, and loss-run issues.",
                "Reconciliation — cross-document provenance when ACORD, inspection, and loss runs disagree.",
                "Pricing build-up — indicated premium with modifier transparency.",
                "Human checkpoints — licensed UW review gates before bind.",
            ],
            "Screenshot: Submission journey drawer for Pacific Coast Distributors",
        ),
        (
            "04_queue",
            "Prioritized submission queue",
            "The Queue prioritizes work so underwriters focus on the highest-value submissions first.",
            [
                "Triage scoring — hot, warm, cold, and no-fit priority bands.",
                "Insured context — name, estimated premium, and bundle ID on every row.",
                "Journey strip — mini pipeline indicator without opening the full drawer.",
                "Click-through — jump directly into the submission journey for any queued item.",
            ],
            "Screenshot: Submission queue with triage filters and journey indicators",
        ),
    ]

    for key, heading, what_it_is, bullets, caption in feature_slides:
        if key in shots and shots[key].exists():
            add_feature_slide(heading, what_it_is, bullets, shots[key], caption)

    add_text_slide(
        "Built for production carriers",
        "Rytera is architected for enterprise deployment, not just demos.",
        [
            "Integration gateway — REST adapters for oracles, policy admin, CRM, and ops systems.",
            "Auto / live / simulated modes — demo without credentials; go live when keys are configured.",
            "Role-based access — viewer, underwriter, licensed UW, admin, and CUO roles.",
            "Encrypted audit bundles — regulatory export with SHA-256 manifest for examiners.",
            "Model registry & override analytics — governance for AI-assisted underwriting.",
        ],
    )

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _accent_bar(slide)
    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.8), Inches(3))
    tf = box.text_frame
    tf.paragraphs[0].text = "See the full pipeline — not just the decision"
    tf.paragraphs[0].font.size = Pt(34)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for line in [
        "Built for carriers, MGAs, banks, and credit unions",
        "Early product demo available",
        "rytera.ai",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = BRAND_LIGHT if "rytera" in line else MUTED
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(14)
    _footer(slide)
    _slide_number(slide)

    prs.save(PPT_PATH)
    return PPT_PATH


def load_existing_shots() -> dict[str, Path]:
    shots: dict[str, Path] = {}
    for key in ("01_overview", "02_insurance", "03_submission_journey", "04_queue"):
        path = OUT_DIR / f"{key}.png"
        if path.exists():
            shots[key] = path
    return shots


def main() -> int:
    refresh = "--refresh" in sys.argv
    shots = {} if refresh else load_existing_shots()
    if not shots:
        try:
            token = api_login()
        except Exception as exc:
            existing = load_existing_shots()
            if existing:
                print(f"Server unavailable, using cached screenshots ({exc})")
                shots = existing
            else:
                print(f"Login failed — is the server running on :8002? {exc}", file=sys.stderr)
                return 1
        else:
            print("Capturing dashboard screenshots…")
            shots = capture_screenshots(token)
    else:
        print(f"Using {len(shots)} existing screenshots (pass --refresh to recapture)")

    print("Building PowerPoint…")
    path = build_ppt(shots)
    print(f"Done: {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
